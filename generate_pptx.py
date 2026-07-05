import sys
import os
import subprocess

# 自動安裝 python-pptx 庫
try:
    import pptx
except ImportError:
    print("未偵測到 python-pptx，正在嘗試自動安裝...")
    try:
        subprocess.run([sys.executable, "-m", "pip", "install", "python-pptx"], check=True)
        import pptx
        print("python-pptx 安裝成功！")
    except Exception as e:
        print(f"安裝 python-pptx 失敗: {e}。請手動運行 'pip install python-pptx'")
        sys.exit(1)

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def create_presentation():
    # 建立簡報物件
    prs = Presentation()
    
    # 設定投影片比例為 16:9 寬螢幕
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # 顏色定義
    COLOR_BG = RGBColor(10, 14, 23)        # #0A0E17 (深背景)
    COLOR_CARD = RGBColor(18, 24, 38)      # #121824 (卡片底色)
    COLOR_TEXT_MAIN = RGBColor(243, 244, 246) # #F3F4F6 (亮白主字)
    COLOR_TEXT_MUTED = RGBColor(156, 163, 175) # #9CA3AF (灰色副字)
    
    COLOR_CYAN = RGBColor(0, 242, 254)     # 青色
    COLOR_PURPLE = RGBColor(139, 92, 246)   # 紫色
    COLOR_AMBER = RGBColor(245, 158, 11)    # 琥珀橙
    COLOR_ROSE = RGBColor(244, 63, 94)      # 玫瑰紅
    COLOR_EMERALD = RGBColor(16, 185, 129)  # 翡翠綠
    COLOR_BORDER = RGBColor(40, 50, 70)     # 卡片框線色

    # 投影片底色與主標題輔助函式
    def set_dark_background(slide):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = COLOR_BG

    def add_slide_header(slide, tag_text, title_text):
        # 增加大分類標籤 (Tag)
        if tag_text:
            tag_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.4))
            tf_tag = tag_box.text_frame
            tf_tag.word_wrap = True
            p_tag = tf_tag.paragraphs[0]
            p_tag.text = tag_text.upper()
            p_tag.font.name = "Arial"
            p_tag.font.size = Pt(11)
            p_tag.font.bold = True
            p_tag.font.color.rgb = COLOR_CYAN
        
        # 增加主標題 (Title)
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.7), Inches(11.7), Inches(0.8))
        tf_title = title_box.text_frame
        tf_title.word_wrap = True
        p_title = tf_title.paragraphs[0]
        p_title.text = title_text
        p_title.font.name = "Microsoft JhengHei"
        p_title.font.size = Pt(28)
        p_title.font.bold = True
        p_title.font.color.rgb = COLOR_TEXT_MAIN

    def draw_card_shape(slide, left, top, width, height, border_color=COLOR_BORDER):
        # 畫卡片底色與框線
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = COLOR_CARD
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1.5)
        return shape

    # ==========================================
    # SLIDE 1: Cover Page
    # ==========================================
    slide_layout = prs.slide_layouts[6] # 空白版面
    slide1 = prs.slides.add_slide(slide_layout)
    set_dark_background(slide1)

    # 裝飾用主視覺圓角矩形
    logo_shape = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.866), Inches(1.2), Inches(1.6), Inches(1.6))
    logo_shape.fill.solid()
    logo_shape.fill.fore_color.rgb = COLOR_CARD
    logo_shape.line.color.rgb = COLOR_CYAN
    logo_shape.line.width = Pt(3)
    tf_logo = logo_shape.text_frame
    tf_logo.word_wrap = True
    p_logo = tf_logo.paragraphs[0]
    p_logo.text = "07-03"
    p_logo.alignment = PP_ALIGN.CENTER
    p_logo.font.name = "Arial"
    p_logo.font.size = Pt(26)
    p_logo.font.bold = True
    p_logo.font.color.rgb = COLOR_CYAN

    # 簡報主標題
    title_box1 = slide1.shapes.add_textbox(Inches(0.5), Inches(3.2), Inches(12.33), Inches(1.2))
    tf1 = title_box1.text_frame
    tf1.word_wrap = True
    p1 = tf1.paragraphs[0]
    p1.text = "分層教學評量系統與 AI 應用"
    p1.alignment = PP_ALIGN.CENTER
    p1.font.name = "Microsoft JhengHei"
    p1.font.size = Pt(40)
    p1.font.bold = True
    p1.font.color.rgb = COLOR_TEXT_MAIN

    # 簡報副標題
    subtitle_box1 = slide1.shapes.add_textbox(Inches(1.0), Inches(4.5), Inches(11.33), Inches(0.8))
    tf_sub = subtitle_box1.text_frame
    tf_sub.word_wrap = True
    p_sub = tf_sub.paragraphs[0]
    p_sub.text = "結合分層結構、RAG 知識防線與即時分析管線的教學評量閉環解決方案"
    p_sub.alignment = PP_ALIGN.CENTER
    p_sub.font.name = "Microsoft JhengHei"
    p_sub.font.size = Pt(18)
    p_sub.font.color.rgb = COLOR_TEXT_MUTED

    # 底部註記
    meta_box = slide1.shapes.add_textbox(Inches(1.0), Inches(5.8), Inches(11.33), Inches(0.8))
    tf_meta = meta_box.text_frame
    tf_meta.word_wrap = True
    p_meta = tf_meta.paragraphs[0]
    p_meta.text = "📅 日期時間：2026-07-03 18:22:27  |  📍 地點：[插入地點]  |  👤 講師：[插入姓名]"
    p_meta.alignment = PP_ALIGN.CENTER
    p_meta.font.name = "Microsoft JhengHei"
    p_meta.font.size = Pt(13)
    p_meta.font.color.rgb = COLOR_TEXT_MUTED

    # ==========================================
    # SLIDE 2: Abstract & Focus
    # ==========================================
    slide2 = prs.slides.add_slide(slide_layout)
    set_dark_background(slide2)
    add_slide_header(slide2, "Lecture Abstract", "講座摘要與核心目標")

    # 左側：講座摘要
    draw_card_shape(slide2, Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.8))
    tb_left = slide2.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(5.2), Inches(4.4))
    tf_left = tb_left.text_frame
    tf_left.word_wrap = True
    
    p_sum_title = tf_left.paragraphs[0]
    p_sum_title.text = "📝 講座摘要說明"
    p_sum_title.font.name = "Microsoft JhengHei"
    p_sum_title.font.size = Pt(18)
    p_sum_title.font.bold = True
    p_sum_title.font.color.rgb = COLOR_CYAN
    p_sum_title.space_after = Pt(12)

    p_sum_desc = tf_left.add_paragraph()
    p_sum_desc.text = "本次討論聚焦於建立一套以「分層結構」為核心的教學與評量系統，結合測驗診斷、學習歷程匯出、個別弱點辨識與 AI 輔助回饋，以支援高爾夫動作要點與骨科解剖學等教學需求。\n\n核心重點在於引進 RAG 限制 AI 回覆以避免幻覺、優化評量管線縮短時差、精簡評語為50字、並對報告加註信賴星等，兼顧師生雙向需求。"
    p_sum_desc.font.name = "Microsoft JhengHei"
    p_sum_desc.font.size = Pt(13)
    p_sum_desc.font.color.rgb = COLOR_TEXT_MUTED
    p_sum_desc.line_spacing = 1.3

    # 右側：三大核心挑戰與解法
    draw_card_shape(slide2, Inches(6.9), Inches(1.8), Inches(5.6), Inches(4.8), border_color=COLOR_AMBER)
    tb_right = slide2.shapes.add_textbox(Inches(7.1), Inches(2.0), Inches(5.2), Inches(4.4))
    tf_right = tb_right.text_frame
    tf_right.word_wrap = True
    
    p_ch_title = tf_right.paragraphs[0]
    p_ch_title.text = "🎯 核心聚焦目標"
    p_ch_title.font.name = "Microsoft JhengHei"
    p_ch_title.font.size = Pt(18)
    p_ch_title.font.bold = True
    p_ch_title.font.color.rgb = COLOR_AMBER
    p_ch_title.space_after = Pt(12)

    challenges = [
        ("1. RAG 知識邊界防線", "排除 AI 生成幻覺，回答嚴格限制在檢索到的專業文獻內。"),
        ("2. 當日分析與決策 SOP", "將傳統長達兩個月的評量滯後，壓縮在當日完成數據對齊。"),
        ("3. 50字回饋與信賴評定", "將發散臨床評語壓縮在 50 字以內，並標註 1-5 星信賴度。")
    ]
    for ch_t, ch_d in challenges:
        p_ct = tf_right.add_paragraph()
        p_ct.text = ch_t
        p_ct.font.name = "Microsoft JhengHei"
        p_ct.font.size = Pt(14)
        p_ct.font.bold = True
        p_ct.font.color.rgb = COLOR_TEXT_MAIN
        
        p_cd = tf_right.add_paragraph()
        p_cd.text = ch_d
        p_cd.font.name = "Microsoft JhengHei"
        p_cd.font.size = Pt(12)
        p_cd.font.color.rgb = COLOR_TEXT_MUTED
        p_cd.space_after = Pt(10)

    # ==========================================
    # SLIDE 3: 4-Layer Standard
    # ==========================================
    slide3 = prs.slides.add_slide(slide_layout)
    set_dark_background(slide3)
    add_slide_header(slide3, "System Architecture", "系統分層與結構設計 (Layer 1 - 4)")

    layers = [
        ("Layer 1: Subjects", "學科核心與知識節點", "定義各學科核心大綱內容（高爾夫動作要點、骨科解剖學知識定義），作為整個系統與 RAG 檢索的知識基石。"),
        ("Layer 2: Quizzes", "測驗題庫與作答紀錄", "包含量化（選擇題、結構化題型）與質性（開放簡答題），負責格式化收集學員現場及線上的作答歷程資料。"),
        ("Layer 3: History", "歷程百分比與分群", "呈現學員學習歷程追蹤，在該學科內判斷學生屬於哪個同儕分群（K-Means 模型）、哪個學習階段，並客製化 PR 值。"),
        ("Layer 4: Reports", "報告輸出與 AI 稽核", "產出學生個人化診斷與教師班級報告（PDF），搭載去識別化敏感資料安全欄位、50 字評語精簡與信賴評分。")
    ]

    col_width = Inches(2.7)
    spacing = Inches(0.3)
    start_left = Inches(0.8)

    for idx, (l_title, l_sub, l_desc) in enumerate(layers):
        left_pos = start_left + idx * (col_width + spacing)
        top_pos = Inches(1.8)
        height_pos = Inches(5.0)
        
        draw_card_shape(slide3, left_pos, top_pos, col_width, height_pos, border_color=COLOR_CYAN)
        
        num_box = slide3.shapes.add_textbox(left_pos + Inches(0.1), top_pos + Inches(0.2), col_width - Inches(0.2), Inches(0.5))
        p_num = num_box.text_frame.paragraphs[0]
        p_num.text = f"LAYER {idx + 1}"
        p_num.alignment = PP_ALIGN.CENTER
        p_num.font.name = "Arial"
        p_num.font.size = Pt(12)
        p_num.font.bold = True
        p_num.font.color.rgb = COLOR_CYAN
        
        tb = slide3.shapes.add_textbox(left_pos + Inches(0.15), top_pos + Inches(0.8), col_width - Inches(0.3), height_pos - Inches(1.0))
        tf = tb.text_frame
        tf.word_wrap = True
        
        p_title = tf.paragraphs[0]
        p_title.text = l_title.split(":")[1].strip()
        p_title.alignment = PP_ALIGN.CENTER
        p_title.font.name = "Microsoft JhengHei"
        p_title.font.size = Pt(18)
        p_title.font.bold = True
        p_title.font.color.rgb = COLOR_TEXT_MAIN
        
        p_sub = tf.add_paragraph()
        p_sub.text = l_sub
        p_sub.alignment = PP_ALIGN.CENTER
        p_sub.font.name = "Microsoft JhengHei"
        p_sub.font.size = Pt(12)
        p_sub.font.color.rgb = COLOR_CYAN
        p_sub.space_after = Pt(14)
        
        p_desc = tf.add_paragraph()
        p_desc.text = l_desc
        p_desc.font.name = "Microsoft JhengHei"
        p_desc.font.size = Pt(12)
        p_desc.font.color.rgb = COLOR_TEXT_MUTED
        p_desc.line_spacing = 1.3

    # ==========================================
    # SLIDE 4: AI & Risk Control
    # ==========================================
    slide4 = prs.slides.add_slide(slide_layout)
    set_dark_background(slide4)
    add_slide_header(slide4, "AI & Risk Control", "AI 應用與風險控管機制")

    draw_card_shape(slide4, Inches(0.8), Inches(1.8), Inches(6.0), Inches(5.0))
    tb_left = slide4.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(5.6), Inches(4.6))
    tf_left = tb_left.text_frame
    tf_left.word_wrap = True

    p_l1 = tf_left.paragraphs[0]
    p_l1.text = "1. RAG 範圍限制與專門聚焦"
    p_l1.font.name = "Microsoft JhengHei"
    p_l1.font.size = Pt(16)
    p_l1.font.bold = True
    p_l1.font.color.rgb = COLOR_CYAN
    
    p_l1_sub = tf_left.add_paragraph()
    p_l1_sub.text = "AI 僅在指定專業內容資料庫內檢索，拒絕回答範疇外問題。讓 AI 回覆高度聚焦於高爾夫動作要點與骨科解剖學考照重點，杜絕 AI 幻覺風險。"
    p_l1_sub.font.name = "Microsoft JhengHei"
    p_l1_sub.font.size = Pt(12)
    p_l1_sub.font.color.rgb = COLOR_TEXT_MUTED
    p_l1_sub.space_after = Pt(14)

    p_l2 = tf_left.add_paragraph()
    p_l2.text = "2. 文字摘要化與信賴星等"
    p_l2.font.name = "Microsoft JhengHei"
    p_l2.font.size = Pt(16)
    p_l2.font.bold = True
    p_l2.font.color.rgb = COLOR_CYAN
    
    p_l2_sub = tf_left.add_paragraph()
    p_l2_sub.text = "將發散的評語，透過 AI 統整為 50 字以內的精簡摘要。以專業 Prompt 生成學員不足面向，並標註 1-5 星信賴等級，以便教學端審閱採納。"
    p_l2_sub.font.name = "Microsoft JhengHei"
    p_l2_sub.font.size = Pt(12)
    p_l2_sub.font.color.rgb = COLOR_TEXT_MUTED
    p_l2_sub.space_after = Pt(14)

    p_l3 = tf_left.add_paragraph()
    p_l3.text = "3. 個別化強化建議"
    p_l3.font.name = "Microsoft JhengHei"
    p_l3.font.size = Pt(16)
    p_l3.font.bold = True
    p_l3.font.color.rgb = COLOR_CYAN
    
    p_l3_sub = tf_left.add_paragraph()
    p_l3_sub.text = "針對弱點對症下藥（例如提示需強化下肢骨頭結構位置），由 AI 或規則引擎提供客製化建議，維持內容可靠性。"
    p_l3_sub.font.name = "Microsoft JhengHei"
    p_l3_sub.font.size = Pt(12)
    p_l3_sub.font.color.rgb = COLOR_TEXT_MUTED

    # 右側範例
    draw_card_shape(slide4, Inches(7.3), Inches(1.8), Inches(5.2), Inches(5.0), border_color=COLOR_EMERALD)
    tb_right = slide4.shapes.add_textbox(Inches(7.5), Inches(2.0), Inches(4.8), Inches(4.6))
    tf_right = tb_right.text_frame
    tf_right.word_wrap = True

    p_r1 = tf_right.paragraphs[0]
    p_r1.text = "💡 RAG 檢索與信賴度評分範例"
    p_r1.font.name = "Microsoft JhengHei"
    p_r1.font.size = Pt(16)
    p_r1.font.bold = True
    p_r1.font.color.rgb = COLOR_EMERALD
    p_r1.space_after = Pt(16)

    p_r2 = tf_right.add_paragraph()
    p_r2.text = "▼ 查詢「斜角肌壓迫與尺神經」"
    p_r2.font.bold = True
    p_r2.font.size = Pt(13)
    p_r2.font.color.rgb = COLOR_TEXT_MAIN
    
    p_r2_sub = tf_right.add_paragraph()
    p_r2_sub.text = "→ AI 50字摘要: 「前中斜角肌與第一肋骨構成隘口壓迫臂神經叢下幹，致第四五指麻木。建議強化斜角肌伸展。」(星級：⭐⭐⭐⭐⭐)"
    p_r2_sub.font.size = Pt(12)
    p_r2_sub.font.color.rgb = COLOR_TEXT_MUTED
    p_r2_sub.space_after = Pt(16)

    p_r3 = tf_right.add_paragraph()
    p_r3.text = "▼ 查詢「高爾夫球動作要點」"
    p_r3.font.bold = True
    p_r3.font.size = Pt(13)
    p_r3.font.color.rgb = COLOR_TEXT_MAIN
    
    p_r3_sub = tf_right.add_paragraph()
    p_r3_sub.text = "→ AI 50字摘要: 「下桿動力鏈應由骨盆率先啟動、重心移向左足。肩膀提早轉動會導致過頂切擊與力量流失。」(星級：⭐⭐⭐⭐)"
    p_r3_sub.font.size = Pt(12)
    p_r3_sub.font.color.rgb = COLOR_TEXT_MUTED

    # ==========================================
    # SLIDE 5: Pipeline & SOP
    # ==========================================
    slide5 = prs.slides.add_slide(slide_layout)
    set_dark_background(slide5)
    add_slide_header(slide5, "Optimization & SOP", "教學評量流程優化與時效")

    # 左側：痛點
    draw_card_shape(slide5, Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.8), border_color=COLOR_ROSE)
    tb_pain = slide5.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(5.2), Inches(4.4))
    tf_pain = tb_pain.text_frame
    tf_pain.word_wrap = True
    
    p_pt = tf_pain.paragraphs[0]
    p_pt.text = "🔴 傳統流程痛點：大量評量與滯後決策"
    p_pt.font.name = "Microsoft JhengHei"
    p_pt.font.size = Pt(18)
    p_pt.font.bold = True
    p_pt.font.color.rgb = COLOR_ROSE
    p_pt.space_after = Pt(16)

    p_pd1 = tf_pain.add_paragraph()
    p_pd1.text = "● 大量手工統計"
    p_pd1.font.bold = True
    p_pd1.font.size = Pt(14)
    p_pd1.font.color.rgb = COLOR_TEXT_MAIN
    
    p_pd1_sub = tf_pain.add_paragraph()
    p_pd1_sub.text = "老師需對大量量化筆試與現場術科評量進行評閱與分析，數據分散不易整合。"
    p_pd1_sub.font.size = Pt(12)
    p_pd1_sub.font.color.rgb = COLOR_TEXT_MUTED
    p_pd1_sub.space_after = Pt(12)

    p_pd2 = tf_pain.add_paragraph()
    p_pd2.text = "● 回饋嚴重滯後 (時差高達兩個月)"
    p_pd2.font.bold = True
    p_pd2.font.size = Pt(14)
    p_pd2.font.color.rgb = COLOR_TEXT_MAIN
    
    p_pd2_sub = tf_pain.add_paragraph()
    p_pd2_sub.text = "學員被評估後資料仍需再整理分析，往往延宕至約兩個月後召開評議會時才得到回饋，顯著影響教學即時性。"
    p_pd2_sub.font.size = Pt(12)
    p_pd2_sub.font.color.rgb = COLOR_TEXT_MUTED

    # 右側：解法
    draw_card_shape(slide5, Inches(6.9), Inches(1.8), Inches(5.6), Inches(4.8), border_color=COLOR_EMERALD)
    tb_sol = slide5.shapes.add_textbox(Inches(7.1), Inches(2.0), Inches(5.2), Inches(4.4))
    tf_sol = tb_sol.text_frame
    tf_sol.word_wrap = True
    
    p_st = tf_sol.paragraphs[0]
    p_st.text = "🟢 新型即時資料管線與會議 SOP"
    p_st.font.name = "Microsoft JhengHei"
    p_st.font.size = Pt(18)
    p_st.font.bold = True
    p_st.font.color.rgb = COLOR_EMERALD
    p_st.space_after = Pt(16)

    p_sd1 = tf_sol.add_paragraph()
    p_sd1.text = "● 當日收集、當日分析、當日會議"
    p_sd1.font.bold = True
    p_sd1.font.size = Pt(14)
    p_sd1.font.color.rgb = COLOR_TEXT_MAIN
    
    p_sd1_sub = tf_sol.add_paragraph()
    p_sd1_sub.text = "於當日收到作答資料即刻自動分析，並召開對齊會議，立即指出學生當週弱點與未來強化方向。"
    p_sd1_sub.font.size = Pt(12)
    p_sd1_sub.font.color.rgb = COLOR_TEXT_MUTED
    p_sd1_sub.space_after = Pt(12)

    p_sd2 = tf_sol.add_paragraph()
    p_sd2.text = "● 教師與學生資訊同步"
    p_sd2.font.bold = True
    p_sd2.font.size = Pt(14)
    p_sd2.font.color.rgb = COLOR_TEXT_MAIN
    
    p_sd2_sub = tf_sol.add_paragraph()
    p_sd2_sub.text = "分析結果與 AI 診斷對師生雙方同時可見，以便於當日即可同步修正學習與教學策略，完成快速改善閉環。"
    p_sd2_sub.font.size = Pt(12)
    p_sd2_sub.font.color.rgb = COLOR_TEXT_MUTED

    # ==========================================
    # SLIDE 6: Scenarios
    # ==========================================
    slide6 = prs.slides.add_slide(slide_layout)
    set_dark_background(slide6)
    add_slide_header(slide6, "Scenarios", "學科場景與應用示例")

    # 左欄高爾夫
    draw_card_shape(slide6, Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.8), border_color=COLOR_CYAN)
    tb_golf = slide6.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(5.2), Inches(4.4))
    tf_golf = tb_golf.text_frame
    tf_golf.word_wrap = True

    p_gt = tf_golf.paragraphs[0]
    p_gt.text = "🏌️ 高爾夫球教學場景"
    p_gt.font.name = "Microsoft JhengHei"
    p_gt.font.size = Pt(18)
    p_gt.font.bold = True
    p_gt.font.color.rgb = COLOR_CYAN
    p_gt.space_after = Pt(14)

    golf_items = [
        ("動作力學大綱定義", "於 Layer 1 定義下桿動力鏈、手腕釋放時機、桿面控制等核心知識節點。"),
        ("測驗與問題診斷", "於 Layer 2 收集作答，自動診斷學員肩膀搶跑（過頂切擊）或手腕釋放過早（鑄型）等動作異常。"),
        ("動作矯正與 PDF 報告", "提供分段重啟動作矯正建議，並匯出包含動作痛點與弱點診斷的個人化 PDF 報告。")
    ]
    for item_t, item_d in golf_items:
        p_it = tf_golf.add_paragraph()
        p_it.text = f"● {item_t}"
        p_it.font.name = "Microsoft JhengHei"
        p_it.font.size = Pt(13)
        p_it.font.bold = True
        p_it.font.color.rgb = COLOR_TEXT_MAIN
        
        p_id = tf_golf.add_paragraph()
        p_id.text = item_d
        p_id.font.name = "Microsoft JhengHei"
        p_id.font.size = Pt(11.5)
        p_id.font.color.rgb = COLOR_TEXT_MUTED
        p_id.space_after = Pt(8)

    # 右欄解剖
    draw_card_shape(slide6, Inches(6.9), Inches(1.8), Inches(5.6), Inches(4.8), border_color=COLOR_PURPLE)
    tb_anat = slide6.shapes.add_textbox(Inches(7.1), Inches(2.0), Inches(5.2), Inches(4.4))
    tf_anat = tb_anat.text_frame
    tf_anat.word_wrap = True

    p_at = tf_anat.paragraphs[0]
    p_at.text = "🦴 骨科解剖學考照場景"
    p_at.font.name = "Microsoft JhengHei"
    p_at.font.size = Pt(18)
    p_at.font.bold = True
    p_at.font.color.rgb = COLOR_PURPLE
    p_at.space_after = Pt(14)

    anat_items = [
        ("大綱與考照痛點", "解決題型、題目長度或空間解剖結構概念不清（如胸廓出口斜角肌隘口、前十字韌帶前拉試驗相對位移）導致的卡關。"),
        ("診斷與定位流程", "測驗診斷 $\rightarrow$ 歷程分群/分階段與同儕百分比 (PR) 定位 $\rightarrow$ 輸出個別弱點與解剖位置建議。"),
        ("臨床物理檢查整合", "整合理學檢查（如脛骨前向位移 ACL 檢測）與解剖起止點定義，系統化落地。")
    ]
    for item_t, item_d in anat_items:
        p_it = tf_anat.add_paragraph()
        p_it.text = f"● {item_t}"
        p_it.font.name = "Microsoft JhengHei"
        p_it.font.size = Pt(13)
        p_it.font.bold = True
        p_it.font.color.rgb = COLOR_TEXT_MAIN
        
        p_id = tf_anat.add_paragraph()
        p_id.text = item_d
        p_id.font.name = "Microsoft JhengHei"
        p_id.font.size = Pt(11.5)
        p_id.font.color.rgb = COLOR_TEXT_MUTED
        p_id.space_after = Pt(8)

    # ==========================================
    # SLIDE 7: Users & Feedback
    # ==========================================
    slide7 = prs.slides.add_slide(slide_layout)
    set_dark_background(slide7)
    add_slide_header(slide7, "Users & Feedback", "使用者存取、自主學習與回饋機制")

    # 左欄使用者
    draw_card_shape(slide7, Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.8))
    tb_u = slide7.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(5.2), Inches(4.4))
    tf_u = tb_u.text_frame
    tf_u.word_wrap = True

    p_ut = tf_u.paragraphs[0]
    p_ut.text = "👥 使用者與存取設計"
    p_ut.font.name = "Microsoft JhengHei"
    p_ut.font.size = Pt(18)
    p_ut.font.bold = True
    p_ut.font.color.rgb = COLOR_CYAN
    p_ut.space_after = Pt(14)

    user_points = [
        ("以學生為主體、兼顧教師端", "使用者以實習學員為主，亦包含老師端檢視分析與教學微調決策。特別適用於內部實習連動之學生。"),
        ("免登入示範環境 (No-Login Demo)", "現階段 Demo 示範不需要學生個別帳號即可觀看樣例，方便先行了解系統功能與資料流。"),
        ("功能擴充流程", "先行觀看 DEMO 確認評量功能與資料管線，後續再擴充權限。")
    ]
    for pt_t, pt_d in user_points:
        p_pt = tf_u.add_paragraph()
        p_pt.text = f"• {pt_t}"
        p_pt.font.name = "Microsoft JhengHei"
        p_pt.font.size = Pt(13)
        p_pt.font.bold = True
        p_pt.font.color.rgb = COLOR_TEXT_MAIN
        
        p_pd = tf_u.add_paragraph()
        p_pd.text = pt_d
        p_pd.font.name = "Microsoft JhengHei"
        p_pd.font.size = Pt(11.5)
        p_pd.font.color.rgb = COLOR_TEXT_MUTED
        p_pd.space_after = Pt(8)

    # 右欄自主學習
    draw_card_shape(slide7, Inches(6.9), Inches(1.8), Inches(5.6), Inches(4.8))
    tb_f = slide7.shapes.add_textbox(Inches(7.1), Inches(2.0), Inches(5.2), Inches(4.4))
    tf_f = tb_f.text_frame
    tf_f.word_wrap = True

    p_ft = tf_f.paragraphs[0]
    p_ft.text = "🔄 自主學習與教學回饋"
    p_ft.font.name = "Microsoft JhengHei"
    p_ft.font.size = Pt(18)
    p_ft.font.bold = True
    p_ft.font.color.rgb = COLOR_PURPLE
    p_ft.space_after = Pt(14)

    feedback_points = [
        ("自主學習支持 (Self-Monitoring)", "透過學習歷程、弱點診斷、PDF 報告匯出，強化學員的自我監測、檢視與學習規劃能力。"),
        ("教學策略動態微調 (Teaching Alignment)", "教師收到 AI 彙整後之診斷報告，可迅速掌握班級同儕分佈與個別弱點，調整下一次教學重點。"),
        ("考核機制系統化落地", "參考並類比臨床 PTY 及 PTY One 考核機制流程，將通過條件與多元評量標準系統化落地。")
    ]
    for pt_t, pt_d in feedback_points:
        p_pt = tf_f.add_paragraph()
        p_pt.text = f"• {pt_t}"
        p_pt.font.name = "Microsoft JhengHei"
        p_pt.font.size = Pt(13)
        p_pt.font.bold = True
        p_pt.font.color.rgb = COLOR_TEXT_MAIN
        
        p_pd = tf_f.add_paragraph()
        p_pd.text = pt_d
        p_pd.font.name = "Microsoft JhengHei"
        p_pd.font.size = Pt(11.5)
        p_pd.font.color.rgb = COLOR_TEXT_MUTED
        p_pd.space_after = Pt(8)

    # ==========================================
    # SLIDE 8: Tasks
    # ==========================================
    slide8 = prs.slides.add_slide(slide_layout)
    set_dark_background(slide8)
    add_slide_header(slide8, "Action Items & Tasks", "會議決策與待辦作業清單 (10 項開發任務)")

    # 10項任務兩列排版
    tb_tasks1 = slide8.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.8))
    tf_t1 = tb_tasks1.text_frame
    tf_t1.word_wrap = True
    
    tasks_col1 = [
        ("1. 明確定義系統分層與資料夾結構", "列出第1層（學科內容）、第2層（測驗）、第3層（學習歷程）、第4層（報告）的檔案與欄位標準。"),
        ("2. 設計並實作第二層測驗題庫", "含量化與質性題型，覆蓋高爾夫動作要點與骨科解剖學考照重點。"),
        ("3. 建立 RAG 知識庫與檢索範圍", "整理並上傳權威教學材料，設定限制以避免 AI 臨床幻覺。"),
        ("4. 開發歷程匯出與弱點診斷模組", "支援同儕分群/階段判定與百分比落點呈現。"),
        ("5. 實作 PDF 報告匯出", "提供學生版（聚焦弱點對策）與教師版（聚焦班級教學調整）雙報告版型。")
    ]
    for t_title, t_desc in tasks_col1:
        p_item = tf_t1.add_paragraph()
        p_item.text = f"❑ {t_title}"
        p_item.font.name = "Microsoft JhengHei"
        p_item.font.size = Pt(12)
        p_item.font.bold = True
        p_item.font.color.rgb = COLOR_AMBER
        
        p_desc = tf_t1.add_paragraph()
        p_desc.text = t_desc
        p_desc.font.name = "Microsoft JhengHei"
        p_desc.font.size = Pt(10.5)
        p_desc.font.color.rgb = COLOR_TEXT_MUTED
        p_desc.space_after = Pt(8)

    tb_tasks2 = slide8.shapes.add_textbox(Inches(6.9), Inches(1.8), Inches(5.6), Inches(4.8))
    tf_t2 = tb_tasks2.text_frame
    tf_t2.word_wrap = True
    
    tasks_col2 = [
        ("6. 建立 AI 回饋彙整流程", "使用 Prompt 將發散回饋壓縮至 50 字以內，並標記 1-5 星信賴等級。"),
        ("7. 規劃即時資料管線與會議 SOP", "自 2026-07-03 起建立當日收集、當日分析、當日會議決策對齊機制。"),
        ("8. 準備免登入示範環境 (DEMO)", "提供免密碼免登入的學員資料模擬環境供師生快速評估功能。"),
        ("9. 設計臨床考核機制", "參考臨床 PTY 流程，明確及格(70分)條件、輔導門檻與改進閉環。"),
        ("10. 制定資料隱私與品質控管", "對敏感資料進行去識別化，明確標記資料來源與版本稽核。")
    ]
    for t_title, t_desc in tasks_col2:
        p_item = tf_t2.add_paragraph()
        p_item.text = f"❑ {t_title}"
        p_item.font.name = "Microsoft JhengHei"
        p_item.font.size = Pt(12)
        p_item.font.bold = True
        p_item.font.color.rgb = COLOR_AMBER
        
        p_desc = tf_t2.add_paragraph()
        p_desc.text = t_desc
        p_desc.font.name = "Microsoft JhengHei"
        p_desc.font.size = Pt(10.5)
        p_desc.font.color.rgb = COLOR_TEXT_MUTED
        p_desc.space_after = Pt(8)

    # 存檔
    output_filename = "project_presentation.pptx"
    prs.save(output_filename)
    print(f"簡報檔案生成成功！已存為: {os.path.abspath(output_filename)}")

if __name__ == "__main__":
    create_presentation()
